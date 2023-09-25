from pptx import Presentation
 
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Random.seed"
subtitle.text = 'Help on method seed in module random:' \
                'seed(a=None, version=2) method of random.Random instance' \
                'Initialize internal state from hashable object.' \
                'None or no argument seeds from current time or from an operating ' \
                'system specific randomness source if available.'
 
title_slide_layout1 = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(title_slide_layout)
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "Random.shuffle"
subtitle1.text = 'Help on method shuffle in module random: \
shuffle(x, random=None) method of random.Random instance \
Shuffle list x in place, and return None.\
Optional argument random is a 0-argument function returning a \
random float in [0.0, 1.0); if it is the default None, the \
standard random.random will be used.'
 
title_slide_layout2 = prs.slide_layouts[0]
slide2 = prs.slides.add_slide(title_slide_layout)
title2 = slide2.shapes.title
subtitle2 = slide2.placeholders[1]
title2.text = "Random.gauss"
subtitle2.text = 'Help on method gauss in module random: \
gauss(mu, sigma) method of random.Random instance \
Gaussian distribution. \
mu is the mean, and sigma is the standard deviation.  This is \
slightly faster than the normalvariate() function. \
Not thread-safe without a lock around calls.'
 
title_slide_layout3 = prs.slide_layouts[0]
slide3 = prs.slides.add_slide(title_slide_layout)
title3 = slide3.shapes.title
subtitle3 = slide3.placeholders[1]
title3.text = "Random.random"
subtitle3.text = 'Help on built-in function random: ' \
                 'random (...)method of random.Random instance' \
                 'random() -> x in the interval [0, 1).'
 
title_slide_layout4 = prs.slide_layouts[0]
slide4 = prs.slides.add_slide(title_slide_layout)
title4 = slide4.shapes.title
subtitle4 = slide4.placeholders[1]
title4.text = "Random.randrange"
subtitle4.text = 'Help on method randrange in module random: \
randrange(start, stop=None, step=1, _int=<class '
'>) method of random.Random instance \
Choose a random item from range(start, stop[, step]). \
This fixes the problem with randint() which includes the \
endpoint; in Python this is usually not what you want.'
 
title_slide_layout5 = prs.slide_layouts[0]
slide5 = prs.slides.add_slide(title_slide_layout)
title5 = slide5.shapes.title
subtitle5 = slide5.placeholders[1]
title5.text = "Random.getrandbits"
subtitle5.text = 'Help on built-in function getrandbits: \
getrandbits(...) method of random.Random instance \
getrandbits(k) -> x.  Generates an int with k random bits.'
prs.save('test.pptx')
